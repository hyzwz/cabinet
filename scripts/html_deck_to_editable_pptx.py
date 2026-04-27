#!/usr/bin/env python3
"""Convert Cabinet's HTML deck format into an editable PPTX.

This intentionally does not screenshot slides. Text is emitted as PPTX text
boxes, cards/steps as editable shapes, and referenced images as separate media.
"""

from __future__ import annotations

import argparse
import re
import sys
from pathlib import Path

try:
    from bs4 import BeautifulSoup
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
except Exception as exc:  # pragma: no cover - exercised manually in missing envs
    print(
        "Missing dependency. Install beautifulsoup4 and python-pptx in the Python environment.",
        file=sys.stderr,
    )
    raise exc


SLIDE_W = Inches(13.333333)
SLIDE_H = Inches(7.5)

INK = RGBColor(10, 31, 61)
INK_TINT = RGBColor(21, 42, 74)
PAPER = RGBColor(241, 243, 245)
PAPER_TINT = RGBColor(228, 232, 236)
ACCENT = RGBColor(190, 156, 78)
MUTED_DARK = RGBColor(95, 108, 126)
MUTED_LIGHT = RGBColor(209, 218, 229)
WHITE = RGBColor(255, 255, 255)

SANS_FONT = "Noto Sans CJK SC"
SERIF_FONT = "Noto Serif CJK SC"
MONO_FONT = "Aptos Mono"


def inches(value: float):
    return Inches(value)


def clean_text(text: str) -> str:
    text = text.replace("\xa0", " ")
    text = re.sub(r"[ \t\r\f\v]+", " ", text)
    text = re.sub(r"\n\s*\n+", "\n", text)
    return text.strip()


def node_text(node) -> str:
    if node is None:
        return ""
    return clean_text(node.get_text("\n", strip=True))


def is_dark(section) -> bool:
    classes = section.get("class") or []
    return "dark" in classes


def add_box(slide, x, y, w, h, fill, line=None, radius=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, inches(x), inches(y), inches(w), inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
      shape.line.fill.background()
    else:
      shape.line.color.rgb = line
      shape.line.width = Pt(1)
    return shape


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    size: float,
    color,
    *,
    bold: bool = False,
    font_face: str = SANS_FONT,
    align=PP_ALIGN.LEFT,
    line_spacing: float | None = None,
):
    text = clean_text(text)
    shape = slide.shapes.add_textbox(inches(x), inches(y), inches(w), inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = inches(0.03)
    frame.margin_right = inches(0.03)
    frame.margin_top = inches(0.02)
    frame.margin_bottom = inches(0.02)
    lines = text.splitlines() or [""]
    for idx, line in enumerate(lines):
        para = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        para.alignment = align
        if line_spacing is not None:
            para.line_spacing = line_spacing
        run = para.add_run()
        run.text = line
        run.font.name = font_face
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return shape


def add_background(slide, dark: bool, hero: bool):
    fill = INK if dark else PAPER
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = fill
    if hero:
        add_box(slide, 0.55, 0.45, 12.2, 6.6, INK_TINT if dark else PAPER_TINT, None, True)


def add_chrome(slide, section, idx: int, total: int, dark: bool):
    color = MUTED_LIGHT if dark else MUTED_DARK
    chrome = section.select_one(".chrome")
    chrome_texts = [node_text(child) for child in chrome.find_all(recursive=False)] if chrome else []
    left = chrome_texts[0] if chrome_texts else "Cabinet Deck"
    right = chrome_texts[1] if len(chrome_texts) > 1 else f"{idx:02d} / {total:02d}"
    add_text(slide, left, 0.62, 0.28, 5.3, 0.22, 7.5, color, font_face=MONO_FONT)
    add_text(slide, right, 11.35, 0.28, 1.35, 0.22, 7.5, color, font_face=MONO_FONT, align=PP_ALIGN.RIGHT)

    foot = section.select_one(".foot")
    foot_texts = [node_text(child) for child in foot.find_all(recursive=False)] if foot else []
    if foot_texts:
        add_text(slide, foot_texts[0], 0.62, 7.02, 6.3, 0.24, 7.5, color, font_face=MONO_FONT)
        if len(foot_texts) > 1:
            add_text(slide, foot_texts[1], 10.7, 7.02, 2.0, 0.24, 7.5, color, font_face=MONO_FONT, align=PP_ALIGN.RIGHT)


def find_main_title(section) -> str:
    node = section.select_one("h1, h2, blockquote")
    return node_text(node)


def find_kicker(section) -> str:
    return node_text(section.select_one(".kicker"))


def find_lead(section) -> str:
    return node_text(section.select_one(".lead"))


def add_card(slide, title: str, body: str, x: float, y: float, w: float, h: float, dark: bool, number: str = ""):
    fill = RGBColor(30, 50, 82) if dark else WHITE
    line = RGBColor(86, 107, 139) if dark else RGBColor(207, 213, 222)
    text = WHITE if dark else INK
    muted = MUTED_LIGHT if dark else MUTED_DARK
    add_box(slide, x, y, w, h, fill, line, True)
    if number:
        add_text(slide, number, x + 0.18, y + 0.16, 0.7, 0.34, 19, ACCENT, bold=True, font_face=SERIF_FONT)
        add_text(slide, title, x + 0.9, y + 0.19, w - 1.05, 0.26, 10, muted, font_face=MONO_FONT)
        body_y = y + 0.62
    else:
        add_text(slide, title, x + 0.18, y + 0.18, w - 0.36, 0.34, 15, text, bold=True, font_face=SERIF_FONT)
        body_y = y + 0.62
    add_text(slide, body, x + 0.18, body_y, w - 0.36, h - 0.78, 10.5, text, font_face=SANS_FONT, line_spacing=1.05)


def add_cards(slide, section, dark: bool):
    stat_cards = section.select(".stat-card")
    pillars = section.select(".pillar")
    steps = section.select(".step")

    cards = []
    if stat_cards:
        for card in stat_cards:
            cards.append((
                node_text(card.select_one(".stat-label")),
                node_text(card.select_one(".stat-note")),
                node_text(card.select_one(".stat-nb")),
            ))
    elif pillars:
        for card in pillars:
            cards.append((
                node_text(card.select_one(".t")),
                node_text(card.select_one(".d")),
                "",
            ))
    elif steps:
        for card in steps:
            cards.append((
                node_text(card.select_one(".step-title")),
                node_text(card.select_one(".step-desc")),
                node_text(card.select_one(".step-nb")),
            ))

    if not cards:
        return

    if len(cards) <= 3:
        cols, rows = len(cards), 1
    elif len(cards) <= 4:
        cols, rows = 2, 2
    else:
        cols, rows = 4, 2

    area_x, area_y, area_w, area_h = 0.82, 3.78, 11.7, 2.75
    if steps:
        area_y, area_h = 2.65, 3.7
    card_w = (area_w - 0.22 * (cols - 1)) / cols
    card_h = (area_h - 0.24 * (rows - 1)) / rows
    for index, (title, body, number) in enumerate(cards):
        col = index % cols
        row = index // cols
        x = area_x + col * (card_w + 0.22)
        y = area_y + row * (card_h + 0.24)
        add_card(slide, title, body, x, y, card_w, card_h, dark, number)


def add_image(slide, html_path: Path, img_node, x: float, y: float, w: float, h: float):
    if img_node is None:
        return
    src = img_node.get("src")
    if not src:
        return
    image_path = (html_path.parent / src).resolve()
    if not image_path.exists():
        return
    slide.shapes.add_picture(str(image_path), inches(x), inches(y), width=inches(w), height=inches(h))


def render_hero(slide, section, idx: int, total: int, dark: bool):
    color = WHITE if dark else INK
    muted = MUTED_LIGHT if dark else MUTED_DARK
    title = find_main_title(section)
    kicker = find_kicker(section)
    lead = find_lead(section)
    meta = node_text(section.select_one(".meta-row"))
    if kicker:
        add_text(slide, kicker, 0.92, 1.25, 6.0, 0.3, 10, ACCENT, font_face=MONO_FONT)
    add_text(slide, title, 0.9, 1.75, 10.5, 2.1, 44 if idx != total else 34, color, bold=True, font_face=SERIF_FONT)
    if lead:
        add_text(slide, lead, 0.94, 4.22, 7.2, 0.8, 18, muted, font_face=SERIF_FONT)
    if meta:
        add_text(slide, meta.replace("\n", " "), 0.95, 5.45, 8.9, 0.3, 8.5, ACCENT, font_face=MONO_FONT)


def render_standard(slide, section, html_path: Path, idx: int, dark: bool):
    color = WHITE if dark else INK
    muted = MUTED_LIGHT if dark else MUTED_DARK
    kicker = find_kicker(section)
    title = find_main_title(section)
    lead = find_lead(section)
    has_image = section.select_one("figure img") is not None
    has_steps = section.select_one(".step") is not None

    if kicker:
        add_text(slide, kicker, 0.78, 0.95, 4.6, 0.28, 9.5, ACCENT, font_face=MONO_FONT)
    add_text(slide, title, 0.75, 1.28, 6.6 if has_image else 10.8, 1.24, 31, color, bold=True, font_face=SERIF_FONT)
    if lead:
        add_text(slide, lead, 0.78, 2.58, 6.2 if has_image else 8.5, 0.64, 15, muted, font_face=SERIF_FONT)

    if has_image:
        add_image(slide, html_path, section.select_one("figure img"), 7.35, 1.55, 5.1, 3.65)
        caption = node_text(section.select_one("figcaption"))
        if caption:
            add_text(slide, caption, 7.35, 5.28, 5.1, 0.25, 8, muted, font_face=MONO_FONT, align=PP_ALIGN.RIGHT)

    if has_steps:
        add_text(slide, "试点路线 / Cabinet 支撑", 0.8, 2.25, 5.2, 0.3, 11, muted, font_face=MONO_FONT)
    add_cards(slide, section, dark)


def convert(html_path: Path, output_path: Path):
    soup = BeautifulSoup(html_path.read_text(encoding="utf-8"), "html.parser")
    sections = soup.select("section.slide")
    if not sections:
        raise ValueError(f"No section.slide elements found in {html_path}")

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank = prs.slide_layouts[6]

    for idx, section in enumerate(sections, start=1):
        slide = prs.slides.add_slide(blank)
        dark = is_dark(section)
        hero = "hero" in (section.get("class") or [])
        add_background(slide, dark, hero)
        add_chrome(slide, section, idx, len(sections), dark)
        if hero:
            render_hero(slide, section, idx, len(sections), dark)
        else:
            render_standard(slide, section, html_path, idx, dark)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)


def main() -> int:
    parser = argparse.ArgumentParser(description="Convert Cabinet HTML deck to editable PPTX.")
    parser.add_argument("html", type=Path, help="Path to index.html for the source HTML deck")
    parser.add_argument("output", type=Path, help="Output .pptx path")
    args = parser.parse_args()
    convert(args.html.resolve(), args.output.resolve())
    print(f"Wrote editable PPTX: {args.output}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
